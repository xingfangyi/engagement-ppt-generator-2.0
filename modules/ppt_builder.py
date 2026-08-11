import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from modules.image_builder import (
    build_overview_images,
    build_strengths_image,
    build_bottom10_panel,
    ABB_RED,
    ABB_LILAC,
)
from modules.action_generator import generate_actions


def _remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def _delete_slide(prs, index):
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def _set_text(shape, text):
    try:
        shape.text = text
    except Exception:
        pass


def _replace_picture_by_order(slide, picture_order, image_path):
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if picture_order >= len(pictures):
        return

    pic = pictures[picture_order]
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    _remove_shape(pic)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def _build_action_text(topic, core_question, actions):
    lines = [f'{topic}: "{core_question}"', ""]
    for action in actions:
        lines.append(f"• {action}")
    return "\n".join(lines)


def build_output_ppt(data, template_path):
    tmp_dir = tempfile.mkdtemp()

    prs = Presentation(template_path)

    # Generate images
    overview_left, overview_right = build_overview_images(data, tmp_dir)

    strengths_img = os.path.join(tmp_dir, "strengths.png")
    opportunities_img = os.path.join(tmp_dir, "opportunities.png")
    build_strengths_image(data["strengths"], "Top Strengths", strengths_img, ABB_LILAC)
    build_strengths_image(data["opportunities"], "Top Opportunities", opportunities_img, ABB_RED)

    bottom_left = os.path.join(tmp_dir, "bottom_left.png")
    bottom_right = os.path.join(tmp_dir, "bottom_right.png")
    build_bottom10_panel(data["bottom_10"][:5], "Bottom 10 (1-5)", bottom_left)
    build_bottom10_panel(data["bottom_10"][5:10], "Bottom 10 (6-10)", bottom_right)

    # Slide 1
    slide1 = prs.slides[0]
    if len(slide1.shapes) > 6:
        _set_text(slide1.shapes[6], f"Based on Engagement Survey in Oct 2025, {data['department']}")

    # Slide 2
    slide2 = prs.slides[1]
    _set_text(slide2.shapes[2], f"{data['department']} - Overview")
    _replace_picture_by_order(slide2, 0, overview_left)
    _replace_picture_by_order(slide2, 1, overview_right)

    # Slide 3
    slide3 = prs.slides[2]
    _set_text(slide3.shapes[2], f"{data['department']} - Strengths and Opportunities / compares to Company")
    _replace_picture_by_order(slide3, 0, strengths_img)
    _replace_picture_by_order(slide3, 1, opportunities_img)

    # Slide 4
    slide4 = prs.slides[3]
    _set_text(slide4.shapes[2], f"{data['department']} - Bottom 10 Scores")
    _replace_picture_by_order(slide4, 0, bottom_left)
    _replace_picture_by_order(slide4, 1, bottom_right)

    # Slide 8
    slide8 = prs.slides[7]
    actions = generate_actions(data)

    # subtitle
    _set_text(slide8.shapes[15], f"{data['department']}, Engagement Survey 2025")

    # topic labels
    _set_text(slide8.shapes[9], f"Improvement topic: {actions['left_topic']}")
    _set_text(slide8.shapes[10], f"Improvement topic: {actions['right_topic']}")

    # main text blocks
    left_text = _build_action_text(
        actions["left_topic"],
        actions["left_question"],
        actions["left_actions"]
    )
    right_text = _build_action_text(
        actions["right_topic"],
        actions["right_question"],
        actions["right_actions"]
    )

    _set_text(slide8.shapes[3], left_text)
    _set_text(slide8.shapes[5], right_text)
    _set_text(slide8.shapes[8], f"Professional Leadership Development {data['department']}")

    # Delete slides 5-7
    _delete_slide(prs, 6)
    _delete_slide(prs, 5)
    _delete_slide(prs, 4)

    output_path = os.path.join(tmp_dir, f"{data['department']}_Development_Actions_2026.pptx")
    prs.save(output_path)
    return output_path
