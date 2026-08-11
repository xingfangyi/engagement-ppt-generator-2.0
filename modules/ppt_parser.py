import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _safe_text(shape):
    try:
        return shape.text.strip()
    except Exception:
        return ""


def _extract_number(text, default=0):
    m = re.search(r"-?\d+", str(text))
    return int(m.group()) if m else default


def _split_question(text):
    parts = [p.strip() for p in text.split("|", 1)]
    if len(parts) == 2:
        return parts[0], parts[1]
    return text.strip(), text.strip()


def parse_survey_ppt(uploaded_file):
    prs = Presentation(uploaded_file)

    data = {
        "department": "Unknown Team",
        "survey_month": "Oct 2025",
        "overall_score": 0,
        "company_score": 0,
        "change_vs_last": 0,
        "responded_count": 0,
        "question_count": 0,
        "strengths": [],
        "opportunities": [],
        "score_rows": [],
        "bottom_10": [],
    }

    # Slide 2: department / metadata
    if len(prs.slides) >= 2:
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            txt = _safe_text(shape)

            if "Manager:" in txt:
                lines = [x.strip() for x in txt.splitlines() if x.strip()]
                if len(lines) >= 2:
                    data["department"] = lines[-1]
                elif "|" in txt:
                    data["department"] = txt.split("|")[-1].strip()

            if "Questions" in txt:
                q = re.search(r"(\d+)\s*/\s*(\d+)\s*Questions", txt, re.I)
                if q:
                    data["question_count"] = int(q.group(1))

            if "October 2025" in txt:
                data["survey_month"] = "Oct 2025"

    # Slide 3: overview
    if len(prs.slides) >= 3:
        slide3 = prs.slides[2]
        for shape in slide3.shapes:
            txt = _safe_text(shape)
            name = getattr(shape, "name", "")

            if name == "Score value":
                data["overall_score"] = _extract_number(txt)

            if name == "Comparison value":
                data["company_score"] = _extract_number(txt)

            if name == "Change value":
                data["change_vs_last"] = _extract_number(txt)

            if "Responded in" in txt:
                m = re.search(r"(\d+)\s*\(", txt)
                if m:
                    data["responded_count"] = int(m.group(1))

    # Slide 4: strengths & opportunities
    if len(prs.slides) >= 4:
        slide4 = prs.slides[3]
        for shape in slide4.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and "Report section group" in shape.name:
                try:
                    drivers_group = shape.shapes[0]
                    driver_groups = list(drivers_group.shapes)

                    for i, dg in enumerate(driver_groups):
                        txt = _safe_text(dg.shapes[0])
                        impact = _safe_text(dg.shapes[1]) if len(dg.shapes) > 1 else ""
                        title, statement = _split_question(txt)
                        title = re.sub(r"\d+$", "", title).strip()

                        item = {
                            "title": title,
                            "statement": statement,
                            "impact": impact,
                        }

                        if i < 3:
                            data["strengths"].append(item)
                        else:
                            data["opportunities"].append(item)
                except Exception:
                    pass

    # Slides 5-12: score table
    for slide in prs.slides:
        title_texts = [_safe_text(s) for s in slide.shapes if _safe_text(s)]
        if not any(t.startswith("Scores") for t in title_texts):
            continue

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and "Report section group" in shape.name:
                try:
                    table_body = shape.shapes[1]
                    for row_group in table_body.shapes:
                        score = _extract_number(_safe_text(row_group.shapes[0]), 0)
                        question_text = _safe_text(row_group.shapes[1])
                        vs_company = _safe_text(row_group.shapes[2])
                        change = _safe_text(row_group.shapes[3])
                        impact = _safe_text(row_group.shapes[4])
                        comments = _safe_text(row_group.shapes[5]) if len(row_group.shapes) > 5 else "--"

                        driver, statement = _split_question(question_text)

                        row = {
                            "score": score,
                            "driver": driver.strip(),
                            "statement": statement.strip(),
                            "question": question_text.strip(),
                            "vs_company": vs_company.strip(),
                            "change": change.strip(),
                            "impact": impact.strip(),
                            "comments": comments.strip(),
                        }
                        data["score_rows"].append(row)
                except Exception:
                    pass

    data["score_rows"] = [r for r in data["score_rows"] if r["score"] > 0]
    data["bottom_10"] = sorted(
        data["score_rows"],
        key=lambda x: (x["score"], x["driver"].lower())
    )[:10]

    return data
